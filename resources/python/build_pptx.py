#!/usr/bin/env python3
"""
Genera un .pptx 16:9 brandizzato da una specifica JSON (Schola P27).

Non reinventa il formato OOXML: usa python-pptx. La spec arriva da stdin, il
path di output da argv[1]. Tutto il disegno avviene su slide BLANK con shape
posizionate a mano (margini/spaziatura controllati), così il tema
(colori/font/logo) è applicato davvero — niente placeholder stock 4:3.

Contratto layout CHIUSO: ogni slide dichiara un `layout` tra quelli noti
(cover, process_cards, columns, stat, bullets_clean). Un layout sconosciuto o
con campi mancanti ricade su bullets_clean (mai una slide rotta).

Spec JSON:
{
  "theme": {
    "ink": "0A0A0A", "background": "F4F1EA", "accent": "A6192E",
    "fonts": {"title": {"primary": "...", "fallback": "..."},
              "body":  {"primary": "...", "fallback": "..."}},
    "logo_path": "/abs/path/logo.png"        # opzionale; solo sulla cover
  },
  "slides": [
    {"layout": "cover", "title": "...", "subtitle": "...", "school": "..."},
    {"layout": "process_cards", "title": "...", "steps": [{"title": "...", "text": "..."}]},
    {"layout": "columns", "title": "...", "columns": [{"icon": "X", "title": "...", "text": "..."}]},
    {"layout": "stat", "title": "...", "value": "70%", "label": "..."},
    {"layout": "bullets_clean", "title": "...", "bullets": ["...", "..."], "notes": "..."}
  ]
}
"""
import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Tela 16:9
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.75
CONTENT_W = SLIDE_W - 2 * MARGIN

KNOWN_LAYOUTS = {"cover", "process_cards", "columns", "stat", "bullets_clean"}


def hex_to_rgb(value, fallback=(85, 177, 174)):
    try:
        value = (value or "").lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*fallback)


def mix(rgb_a, rgb_b, t):
    """Interpolazione lineare tra due RGBColor (t=0 -> a, t=1 -> b)."""
    return RGBColor(
        int(rgb_a[0] + (rgb_b[0] - rgb_a[0]) * t),
        int(rgb_a[1] + (rgb_b[1] - rgb_a[1]) * t),
        int(rgb_a[2] + (rgb_b[2] - rgb_a[2]) * t),
    )


class Theme:
    def __init__(self, spec_theme):
        spec_theme = spec_theme or {}
        self.ink = hex_to_rgb(spec_theme.get("ink", "0A0A0A"), (10, 10, 10))
        self.bg = hex_to_rgb(spec_theme.get("background", "F4F1EA"), (244, 241, 234))
        self.accent = hex_to_rgb(spec_theme.get("accent", "A6192E"), (166, 25, 46))
        fonts = spec_theme.get("fonts") or {}
        title = fonts.get("title") or {}
        body = fonts.get("body") or {}
        self.title_font = title.get("primary") or "Georgia"
        self.body_font = body.get("primary") or "Calibri"
        self.logo_path = spec_theme.get("logo_path")
        # toni derivati
        self.muted_on_dark = mix(self.bg, self.ink, 0.30)            # testo 2ndario su cover scura
        self.accent_on_dark = mix(self.accent, RGBColor(255, 255, 255), 0.22)  # accento ravvivato su scuro
        self.surface = mix(self.bg, RGBColor(255, 255, 255), 0.55)   # card su sfondo chiaro
        self.hairline = mix(self.bg, self.ink, 0.18)


# ---------------------------------------------------------------- helpers


def set_bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def add_para(tf, text, font, size, color, bold=False, align=PP_ALIGN.LEFT,
             space_after=6, space_before=0, line=1.0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    run = p.add_run()
    run.text = text if text is not None else ""
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def fill_shape(shape, rgb, line_rgb=None, line_w=0.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if line_rgb is not None:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False


def title_block(slide, theme, text, color=None):
    """Titolo di sezione + filetto accento. Ritorna il top sotto il titolo."""
    color = color or theme.accent
    _, tf = textbox(slide, MARGIN, 0.55, CONTENT_W, 1.0)
    add_para(tf, text or "", theme.title_font, 30, color, bold=True, line=1.05, first=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(1.55),
                                  Inches(2.2), Pt(3))
    fill_shape(rule, theme.accent)
    return 1.9


# ---------------------------------------------------------------- layouts


def render_cover(slide, prs, theme, d):
    # Sfondo scuro = ink del tema (sempre scuro nei 4 temi curati).
    set_bg(slide, theme.ink)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(2.4),
                                 Inches(0.18), Inches(2.4))
    fill_shape(bar, theme.accent)

    _, tf = textbox(slide, MARGIN + 0.45, 2.3, CONTENT_W - 0.45, 3.0)
    add_para(tf, d.get("title", "Lezione"), theme.title_font, 44, theme.bg,
             bold=True, line=1.05, first=True)
    sub = d.get("subtitle")
    if sub:
        add_para(tf, sub, theme.body_font, 20, theme.accent_on_dark, space_before=10)

    school = d.get("school")
    if school:
        _, sf = textbox(slide, MARGIN + 0.45, SLIDE_H - 1.2, CONTENT_W - 2.0, 0.6,
                        anchor=MSO_ANCHOR.BOTTOM)
        add_para(sf, school, theme.body_font, 14, theme.muted_on_dark, first=True)

    _logo(slide, prs, theme)


def render_bullets_clean(slide, prs, theme, d):
    set_bg(slide, theme.bg)
    top = title_block(slide, theme, d.get("title"))

    bullets = [b for b in (d.get("bullets") or []) if str(b).strip()]
    if not bullets:
        bullets = [d.get("subtitle") or " "]
    # Centrato verticalmente nell'area disponibile: niente grande vuoto in basso.
    _, tf = textbox(slide, MARGIN, top, CONTENT_W, SLIDE_H - top - MARGIN, anchor=MSO_ANCHOR.MIDDLE)
    for i, b in enumerate(bullets):
        add_para(tf, "•  " + str(b), theme.body_font, 19, theme.ink,
                 space_after=14, line=1.12, first=(i == 0))


def render_process_cards(slide, prs, theme, d):
    set_bg(slide, theme.bg)
    top = title_block(slide, theme, d.get("title"))

    steps = [s for s in (d.get("steps") or []) if isinstance(s, dict)][:5]
    n = len(steps)
    if n == 0:
        return render_bullets_clean(slide, prs, theme, d)

    gap = 0.35
    card_w = (CONTENT_W - gap * (n - 1)) / n
    # Card di altezza contenuta e centrate verticalmente (niente vuoto enorme in fondo).
    avail_h = (SLIDE_H - MARGIN) - top
    card_h = min(avail_h, 3.1)
    card_top = top + (avail_h - card_h) / 2

    for i, s in enumerate(steps):
        left = MARGIN + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(card_top),
                                      Inches(card_w), Inches(card_h))
        fill_shape(card, theme.surface, line_rgb=theme.hairline, line_w=1.0)

        dia = 0.55
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(left + 0.22), Inches(card_top + 0.25),
                                      Inches(dia), Inches(dia))
        fill_shape(circ, theme.accent)
        add_para(circ.text_frame, str(i + 1), theme.title_font, 18, theme.bg, bold=True,
                 align=PP_ALIGN.CENTER, first=True)

        _, tf = textbox(slide, left + 0.22, card_top + 1.0, card_w - 0.44, card_h - 1.2)
        add_para(tf, str(s.get("title", "")).strip(), theme.title_font, 15,
                 theme.ink, bold=True, space_after=6, line=1.05, first=True)
        txt = str(s.get("text", "")).strip()
        if txt:
            add_para(tf, txt, theme.body_font, 11.5, theme.ink, line=1.12)

        if i < n - 1:
            af = slide.shapes.add_textbox(Inches(left + card_w - 0.05),
                                          Inches(card_top + card_h / 2 - 0.35),
                                          Inches(gap + 0.1), Inches(0.7))
            add_para(af.text_frame, "›", theme.title_font, 30, theme.accent,
                     bold=True, align=PP_ALIGN.CENTER, first=True)


def render_columns(slide, prs, theme, d):
    set_bg(slide, theme.bg)
    top = title_block(slide, theme, d.get("title"))

    cols = [c for c in (d.get("columns") or []) if isinstance(c, dict)][:3]
    n = len(cols)
    if n < 2:
        return render_bullets_clean(slide, prs, theme, d)

    gap = 0.5
    col_w = (CONTENT_W - gap * (n - 1)) / n
    # Blocco colonne centrato verticalmente nell'area disponibile.
    avail_h = (SLIDE_H - MARGIN) - top
    col_top = top + (avail_h - 2.4) / 2

    for i, c in enumerate(cols):
        left = MARGIN + i * (col_w + gap)
        dia = 0.9
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(left + col_w / 2 - dia / 2), Inches(col_top),
                                      Inches(dia), Inches(dia))
        fill_shape(circ, theme.accent)
        # Iniziale del titolo come "icona": glifo sempre disponibile (le emoji
        # rendono come tofu quando manca il font emoji). Niente caratteri rotti.
        label = (str(c.get("title", "")).strip()[:1] or "•").upper()
        add_para(circ.text_frame, label, theme.title_font, 26, theme.bg, bold=True,
                 align=PP_ALIGN.CENTER, first=True)

        _, tf = textbox(slide, left, col_top + 1.1, col_w, SLIDE_H - col_top - 1.1 - MARGIN)
        add_para(tf, str(c.get("title", "")).strip(), theme.title_font, 17, theme.ink,
                 bold=True, align=PP_ALIGN.CENTER, space_after=8, line=1.05, first=True)
        txt = str(c.get("text", "")).strip()
        if txt:
            add_para(tf, txt, theme.body_font, 13, theme.ink, align=PP_ALIGN.CENTER, line=1.15)


def render_stat(slide, prs, theme, d):
    set_bg(slide, theme.bg)
    title = str(d.get("title", "")).strip()
    if title:
        title_block(slide, theme, title)

    value = str(d.get("value", "")).strip() or "—"
    label = str(d.get("label", "")).strip()

    _, tf = textbox(slide, MARGIN, 2.4, CONTENT_W, 2.6, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, value, theme.title_font, 110, theme.accent, bold=True,
             align=PP_ALIGN.CENTER, line=1.0, first=True)
    if label:
        _, lf = textbox(slide, MARGIN + 1.0, 5.1, CONTENT_W - 2.0, 1.4)
        add_para(lf, label, theme.body_font, 22, theme.ink, align=PP_ALIGN.CENTER,
                 line=1.15, first=True)


RENDERERS = {
    "cover": render_cover,
    "process_cards": render_process_cards,
    "columns": render_columns,
    "stat": render_stat,
    "bullets_clean": render_bullets_clean,
}


def _logo(slide, prs, theme):
    """Inserisce il logo (se presente e leggibile) in alto a destra. Mai bloccante."""
    path = theme.logo_path
    if not path:
        return
    try:
        pic = slide.shapes.add_picture(path, Inches(0), Inches(0), height=Inches(0.62))
        pic.left = int(prs.slide_width - Inches(MARGIN) - pic.width)
        pic.top = Inches(0.5)
    except Exception:
        pass


def main():
    spec = json.load(sys.stdin)
    out_path = sys.argv[1]
    theme = Theme(spec.get("theme"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]  # Blank

    slides = spec.get("slides") or []
    for s in slides:
        if not isinstance(s, dict):
            continue
        layout = s.get("layout") if s.get("layout") in KNOWN_LAYOUTS else "bullets_clean"
        renderer = RENDERERS.get(layout, render_bullets_clean)
        slide = prs.slides.add_slide(blank)
        try:
            renderer(slide, prs, theme, s)
        except Exception:
            # Difesa: l'errore di un renderer non rompe la presentazione.
            render_bullets_clean(slide, prs, theme, s)

        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(out_path)
    print(json.dumps({"ok": True, "slides": len(slides)}))


if __name__ == "__main__":
    main()
